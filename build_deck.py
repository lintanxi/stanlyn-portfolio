#!/usr/bin/env python3
# Build an editable, Keynote-compatible PPTX deck with speaker notes.
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ASSETS = "assets"
FONT = "Helvetica Neue"

def C(h): return RGBColor.from_string(h)
BG=C("F0EFEB"); SURFACE=C("F7F6F2"); INK=C("222222"); MUTED=C("515357")
LINE=C("E1E0D8"); BLUE=C("0A5FD6"); BLUEINK=C("002E99"); BLUEPALE=C("E6ECF9")
LEAF=C("3A9D3A"); RED=C("C0392B"); REDPALE=C("FDF0EE"); GREENPALE=C("E8F5E9")
GREENINK=C("1B5E20"); WHITE=C("FFFFFF"); LT=C("C9D2DA"); LT2=C("8FA0AD"); SKY=C("6EA8F0")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5

def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

def rect(s, l, t, w, h, fill=None, line=None, line_w=1.0, radius=0.06,
         line_left=None, line_left_w=4):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    try:
        if radius: shp.adjustments[0] = radius
    except Exception: pass
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    return shp

def text(s, l, t, w, h, lines, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left","margin_right","margin_top","margin_bottom"):
        setattr(tf, m, 0)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", PP_ALIGN.LEFT)
        p.space_before = Pt(ln.get("sb", 0))
        p.space_after = Pt(ln.get("sa", 0))
        if "ls" in ln: p.line_spacing = ln["ls"]
        size = ln.get("size", 16); color = ln.get("color", MUTED)
        bold = ln.get("bold", False); font = ln.get("font", FONT)
        if "runs" in ln:
            for (txt, rc, rb) in ln["runs"]:
                r = p.add_run(); r.text = txt
                r.font.size = Pt(size); r.font.name = font
                r.font.bold = rb; r.font.color.rgb = rc
        else:
            r = p.add_run(); r.text = ln.get("text", "")
            r.font.size = Pt(size); r.font.name = font
            r.font.bold = bold; r.font.color.rgb = color
    return tb

def in_shape(shp, lines, anchor=MSO_ANCHOR.TOP):
    tf = shp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.16); tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.12); tf.margin_bottom = Inches(0.12)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", PP_ALIGN.LEFT)
        p.space_before = Pt(ln.get("sb", 0)); p.space_after = Pt(ln.get("sa", 0))
        if "ls" in ln: p.line_spacing = ln["ls"]
        size = ln.get("size", 13); color = ln.get("color", MUTED)
        bold = ln.get("bold", False)
        if "runs" in ln:
            for (txt, rc, rb) in ln["runs"]:
                r = p.add_run(); r.text = txt
                r.font.size = Pt(size); r.font.name = FONT
                r.font.bold = rb; r.font.color.rgb = rc
        else:
            r = p.add_run(); r.text = ln.get("text", "")
            r.font.size = Pt(size); r.font.name = FONT
            r.font.bold = bold; r.font.color.rgb = color

def picture(s, path, l, t, w=None, h=None, box_w=None, box_h=None):
    """Fit image inside box (box_w x box_h) preserving aspect, centered."""
    full = os.path.join(ASSETS, path)
    iw, ih = Image.open(full).size
    ar = iw / ih
    if box_w and box_h:
        bw, bh = box_w, box_h
        if bw / bh > ar:   # box wider -> constrain by height
            ph = bh; pw = bh * ar
        else:
            pw = bw; ph = bw / ar
        pl = l + (bw - pw) / 2; pt = t + (bh - ph) / 2
        return s.shapes.add_picture(full, Inches(pl), Inches(pt), Inches(pw), Inches(ph))
    if w: return s.shapes.add_picture(full, Inches(l), Inches(t), width=Inches(w))
    return s.shapes.add_picture(full, Inches(l), Inches(t), height=Inches(h))

def kicker(s, txt, color=BLUE, l=0.78, t=0.62):
    text(s, l, t, 11, 0.4, [{"text": txt.upper(), "size": 12, "bold": True, "color": color}])

def footer(s, label):
    text(s, 0.78, 7.0, 6, 0.3, [{"text": label, "size": 11, "bold": True, "color": MUTED}])

def pagenum(s, n):
    text(s, 11.3, 7.0, 1.5, 0.3, [{"text": n, "size": 11, "bold": True, "color": MUTED, "align": PP_ALIGN.RIGHT}])

def dot_eyebrow(s, txt, l=0.78, t=0.62):
    rect(s, l, t+0.045, 0.09, 0.09, fill=BLUE, radius=0.5)
    text(s, l+0.2, t, 11, 0.4, [{"text": txt.upper(), "size": 12, "bold": True, "color": MUTED}])

# ============================================================ COVER
s = slide(INK)
picture(s, "../favicon-placeholder", 0, 0) if False else None
# leaf logo drawn as circle
logo = rect(s, 6.42, 1.7, 0.5, 0.5, fill=C("2D8841"), radius=0.5)
text(s, 0, 2.55, SW, 0.4, [{"text": "PRODUCT DESIGN · CASE STUDIES", "size": 13, "bold": True,
     "color": LT2, "align": PP_ALIGN.CENTER}])
text(s, 1, 3.0, SW-2, 1.6, [
    {"runs": [("Designing systems, ", WHITE, True), ("shipping services", SKY, True)],
     "size": 44, "align": PP_ALIGN.CENTER, "ls": 1.04}], anchor=MSO_ANCHOR.TOP)
text(s, 2.66, 4.6, SW-5.32, 1.0, [
    {"text": "Two case studies from Walmart Global Tech — one about designing the system "
     "behind the system, one about unifying a service to global scale.",
     "size": 18, "color": LT, "align": PP_ALIGN.CENTER, "ls": 1.4}])
text(s, 0, 5.85, SW, 0.4, [{"text": "Stanlyn Lu · Product Designer", "size": 15, "bold": True,
     "color": WHITE, "align": PP_ALIGN.CENTER}])
notes(s, "Open with the framing line: 'My first case is about designing the system behind the "
       "system; my second is about shipping a unified service to global scale. Together they show "
       "how I move between strategy and execution.' Keep this to ~15 seconds, then advance.")

# ============================================================ FRAMING
s = slide()
dot_eyebrow(s, "How to read these two cases")
text(s, 0.78, 1.15, 11, 0.7, [{"text": "Two cases, chosen to show range", "size": 32, "bold": True, "color": INK}])
c1 = rect(s, 0.78, 2.2, 5.7, 4.2, fill=SURFACE, line=LINE)
rect(s, 0.78, 2.2, 5.7, 0.06, fill=BLUE, radius=0)
in_shape(c1, [
    {"text": "CASE 01 · WARDEN", "size": 12, "bold": True, "color": BLUE, "sa": 8},
    {"text": "0→1 systems & governance", "size": 22, "bold": True, "color": INK, "sa": 10, "ls": 1.05},
    {"text": "I reframed a usability request as a structural problem and designed the infrastructure "
     "that lets operations change the system safely — without going through engineering.",
     "size": 15, "color": MUTED, "sa": 10, "ls": 1.3},
    {"runs": [("Signals: ", INK, True), ("strategic thinking, working solo in ambiguity, AI "
     "judgment in high-stakes workflows.", MUTED, False)], "size": 15, "ls": 1.3},
], anchor=MSO_ANCHOR.TOP)
c2 = rect(s, 6.85, 2.2, 5.7, 4.2, fill=SURFACE, line=LINE)
rect(s, 6.85, 2.2, 5.7, 0.06, fill=LEAF, radius=0)
in_shape(c2, [
    {"text": "CASE 02 · VISIT CHECK-IN", "size": 12, "bold": True, "color": LEAF, "sa": 8},
    {"text": "Service design & measured impact", "size": 22, "bold": True, "color": INK, "sa": 10, "ls": 1.05},
    {"text": "I unified visitor check-in from twelve isolated, inaccessible systems into one "
     "accessible service across every US and India campus.",
     "size": 15, "color": MUTED, "sa": 10, "ls": 1.3},
    {"runs": [("Signals: ", INK, True), ("craft, accessibility rigor, stakeholder alignment, "
     "real outcomes (30% / 3× / ~60%).", MUTED, False)], "size": 15, "ls": 1.3},
], anchor=MSO_ANCHOR.TOP)
footer(s, "Stanlyn Lu"); pagenum(s, "Framing")
notes(s, "Tell the interviewer how to read the pair before diving in: Warden = systems/strategy, "
       "Visit = service/execution + numbers. This primes them to see your range rather than two "
       "similar projects. One decision you can mention: Visit Admin is folded into Case 2 as the "
       "operational backbone, not a separate third case.")

def divider(case, title_runs, sub, accent):
    s = slide(INK)
    text(s, 0, 2.35, SW, 0.4, [{"text": case.upper(), "size": 13, "bold": True, "color": accent,
         "align": PP_ALIGN.CENTER}])
    text(s, 1, 2.85, SW-2, 1.8, [{"runs": title_runs, "size": 48, "align": PP_ALIGN.CENTER, "ls": 1.04}])
    text(s, 2.4, 4.95, SW-4.8, 0.8, [{"text": sub, "size": 18, "color": LT,
         "align": PP_ALIGN.CENTER, "ls": 1.4}])
    return s

# ============================================================ DIVIDER WARDEN
s = divider("Case 01",
    [("Governing the ", WHITE, True), ("system", SKY, True), (" behind the system", WHITE, True)],
    "Walmart's global security operations platform · Sole Product Designer", SKY)
notes(s, "Transition line: 'My first case starts with a reframe — I was asked to improve usability, "
       "and discovered the real problem was governance.' Pause here so the reframe lands before you "
       "show context.")

# ============================================================ W1 title
s = slide()
dot_eyebrow(s, "Governance Design · Security Operations · Design Systems · AI")
text(s, 0.78, 1.1, 8.4, 1.3, [{"runs": [("I was asked to improve usability. The real problem was ",
     INK, True), ("governance.", BLUE, True)], "size": 30, "ls": 1.08}])
text(s, 0.78, 2.7, 5.85, 1.5, [{"text": "How I reframed a usability request as a governance "
     "problem, and designed the infrastructure that lets the people closest to operations change "
     "the system safely.", "size": 17, "color": MUTED, "ls": 1.45}])
# meta strip
meta = [("Role","Sole Product Designer"),("Surfaces","Warden · Mobile · Emergency App"),
        ("Team","~40 · Product + Eng"),("Timeline","Feb 2026 – Present")]
mw = 5.85/4
rect(s, 0.78, 4.5, 5.85, 1.2, fill=SURFACE, line=LINE)
for i,(lab,val) in enumerate(meta):
    if i>0: rect(s, 0.78+i*mw, 4.5, 0.012, 1.2, fill=LINE, radius=0)
    text(s, 0.92+i*mw, 4.66, mw-0.2, 1.0, [
        {"text": lab.upper(), "size": 9, "bold": True, "color": MUTED, "sa": 5},
        {"text": val, "size": 12, "bold": True, "color": INK, "ls": 1.15}])
picture(s, "warden.jpg", 6.95, 1.1, box_w=5.6, box_h=4.6)
rect(s, 6.95, 1.1, 5.6, 4.6, line=LINE)
footer(s, "Warden"); pagenum(s, "01")
notes(s, "One sentence on what Warden is — Walmart's global security operations ecosystem, real "
       "incidents, high stakes — then state the headline reframe. Don't over-explain the org chart; "
       "the meta strip carries the facts. Land the line: 'They asked for usability; I delivered "
       "governance infrastructure.'")

# ============================================================ W2 situation
s = slide()
kicker(s, "The situation I inherited")
text(s, 0.78, 1.05, 11, 0.7, [{"text": "A platform that had outgrown its foundations", "size": 30, "bold": True, "color": INK}])
text(s, 0.78, 1.8, 11.7, 0.9, [{"runs": [("Warden was built fast — growth to ", MUTED, False),
     ("72+ incident types across 11 categories", INK, True), (", more regions and compliance, a "
     "designer transition — left a system never designed to be maintained at this scale. I joined "
     "as the sole designer on a ~40-person team.", MUTED, False)], "size": 16, "ls": 1.4}])
cards = [
    ("Operational Complexity","72+ incident types, 11+ categories, 80+ question categories, a taxonomy still evolving. Regional and compliance rules with no unified logic."),
    ("UX Debt","A codebase mixing Material UI, Ant Design, and custom styling. Inconsistent patterns across surfaces; a prior sub-system only partially adopted."),
    ("Governance Bottleneck","Every questionnaire change needed a spreadsheet, approval meetings, an engineering ticket, and a release cycle. Needs moved faster than the system could adapt."),
    ("Solo Coverage","As the only designer: new feature design, design-system alignment across three layers, and AI evaluation — with no handover from the previous designer."),
]
gx, gy, gw, gh, gap = 0.78, 3.0, 5.78, 1.75, 0.16
for i,(tt,bb) in enumerate(cards):
    cx = gx + (i%2)*(gw+gap); cy = gy + (i//2)*(gh+gap)
    cd = rect(s, cx, cy, gw, gh, fill=SURFACE, line=LINE)
    in_shape(cd, [{"text": tt.upper(), "size": 12, "bold": True, "color": INK, "sa": 6},
                  {"text": bb, "size": 12.5, "color": MUTED, "ls": 1.3}])
footer(s, "Warden"); pagenum(s, "02")
notes(s, "Be honest you walked into this solo with no handover — interviewers want to hear how you "
       "operate under ambiguity. Don't read all four cards; say 'four problems, but they all rhymed,' "
       "and point at Governance Bottleneck as the thread you'll pull next.")

# ============================================================ W3 research
s = slide()
kicker(s, "Stakeholder Research")
text(s, 0.78, 1.05, 11, 0.7, [{"text": "Three different problems. One root cause.", "size": 30, "bold": True, "color": INK}])
text(s, 0.78, 1.8, 11.7, 0.6, [{"text": "I interviewed three stakeholders with different vantage "
     "points — operations, the questionnaire content owner, and the PM who bore the cost of every change.",
     "size": 16, "color": MUTED, "ls": 1.4}])
people = [
    ("Rick","Director, Operations & Technology · GSOC",
     [("“When something changes operationally, we need to ", MUTED, False),
      ("update the questionnaire quickly", BLUEINK, True),
      (" — but by the time a change reaches operators, the situation has moved on.”", MUTED, False)]),
    ("Lizzy","Director, Security Operations · Content Owner",
     [("“I maintain a ", MUTED, False), ("huge spreadsheet of questions", BLUEINK, True),
      (", and every change requires reviews and approvals from multiple leaders before anything goes to engineering.”", MUTED, False)]),
    ("Mike","Product Manager · Warden",
     [("“Every questionnaire change becomes an ", MUTED, False), ("engineering effort", BLUEINK, True),
      (". We need to reduce friction and manual work without creating compliance or data risk.”", MUTED, False)]),
]
cw = 3.78; cgap = 0.18
for i,(nm,ti,runs) in enumerate(people):
    cx = 0.78 + i*(cw+cgap)
    cd = rect(s, cx, 2.85, cw, 3.4, fill=SURFACE, line=LINE)
    in_shape(cd, [{"text": nm, "size": 15, "bold": True, "color": INK, "sa": 2},
                  {"text": ti, "size": 11.5, "color": MUTED, "sa": 12},
                  {"runs": runs, "size": 14, "ls": 1.35}])
footer(s, "Warden"); pagenum(s, "03")
notes(s, "This is your 'I listen before I design' proof. Walk through how three different jobs "
       "pointed at the same structural thing. Quote Lizzy's spreadsheet line verbatim — it's the "
       "most vivid and sets up the next slide.")

# ============================================================ W4 reframe
s = slide()
kicker(s, "Discovery · The reframe")
text(s, 0.78, 1.05, 7.4, 1.0, [{"runs": [("The real problem was workflow ", INK, True),
     ("governance", BLUE, True), (", not the workflow", INK, True)], "size": 29, "bold": True, "ls": 1.05}])
text(s, 0.78, 2.5, 6.4, 1.3, [{"runs": [("A single change followed this path: ", MUTED, False),
     ("spreadsheet → approval meetings → engineering ticket → release cycle.", INK, True),
     (" Weeks, for what was essentially a content edit.", MUTED, False)], "size": 16, "ls": 1.4}])
text(s, 0.78, 3.85, 6.4, 1.3, [{"runs": [("Worse, a simple change had ", MUTED, False),
     ("downstream effects on compliance, dashboards, and notifications", INK, True),
     (" — none visible to the requester. That invisibility was the real problem.", MUTED, False)],
     "size": 16, "ls": 1.4}])
co = rect(s, 0.78, 5.35, 6.4, 1.35, fill=SURFACE, line=LINE)
rect(s, 0.78, 5.35, 0.06, 1.35, fill=BLUE, radius=0)
in_shape(co, [{"text": "CROSS-INTERVIEW INSIGHT", "size": 11, "bold": True, "color": MUTED, "sa": 6},
              {"text": "Workflow ownership was trapped in spreadsheets, approvals, and release cycles. "
               "The people who understood what needed to change had no path to change it.",
               "size": 15, "color": INK, "ls": 1.3}])
picture(s, "warden-pain-pipeline.jpg", 7.6, 1.6, box_w=4.95, box_h=4.5)
footer(s, "Warden"); pagenum(s, "04")
notes(s, "This is the intellectual core of the case — slow down. The point that makes it a DESIGN "
       "problem rather than an ops complaint: one content edit silently cascades into compliance, "
       "dashboards, and notifications. The invisibility is the villain. Don't rush to the solution.")

# ============================================================ W5 solution
s = slide()
kicker(s, "0→1 Design")
text(s, 0.78, 1.05, 11.7, 0.7, [{"text": "Designing the boundary between flexibility and governance",
     "size": 28, "bold": True, "color": INK}])
picture(s, "warden-super-admin.jpg", 0.78, 2.0, box_w=6.9, box_h=4.4)
rect(s, 0.78, 2.0, 6.9, 4.4, line=LINE)
text(s, 8.0, 2.0, 4.55, 0.9, [{"runs": [("The question was ", MUTED, False),
     ("what to allow, not just what to show.", INK, True), (" Phase 1 gave Super Admins four "
     "configurable dimensions:", MUTED, False)], "size": 15, "ls": 1.35}])
dims = [("Visibility","show/hide per question"),("Ordering","drag-to-reorder"),
        ("Requiredness","required vs. optional"),("Country rules","region-based logic")]
for i,(tt,dd) in enumerate(dims):
    ty = 3.55 + i*0.5
    rect(s, 8.0, ty+0.04, 0.16, 0.16, fill=LEAF, radius=0.5)
    text(s, 8.28, ty, 4.3, 0.45, [{"runs": [(tt+" — ", INK, True), (dd, MUTED, False)], "size": 14}])
text(s, 8.0, 5.75, 4.55, 0.9, [{"runs": [("Every change is logged to a persistent ", MUTED, False),
     ("audit trail", INK, True), (" — surfaced in the view, not buried. Auditability was a "
     "first-class requirement.", MUTED, False)], "size": 14, "ls": 1.3}])
footer(s, "Warden"); pagenum(s, "05")
notes(s, "Emphasize the framing 'what to allow, not what to show' — that's the systems-design "
       "insight. Then the audit trail: you surfaced it deliberately because accountability is the "
       "whole point of governance. Four dimensions, but don't recite — point and summarize.")

# ============================================================ W6 what NOT to build
s = slide()
kicker(s, "Design Strategy · The hardest decision")
text(s, 0.78, 1.05, 11.7, 0.7, [{"text": "Earning trust before expanding capability", "size": 30, "bold": True, "color": INK}])
text(s, 0.78, 1.85, 11.7, 0.9, [{"runs": [("The most important decision was what ", MUTED, False),
     ("not", INK, True), (" to build yet. Question authoring was excluded from Phase 1 — not for "
     "technical reasons, but because the taxonomy was mid-consolidation and premature access would "
     "have created real compliance exposure.", MUTED, False)], "size": 16, "ls": 1.4}])
# decision cards
d1 = rect(s, 0.78, 3.05, 5.78, 1.85, fill=SURFACE, line=LINE)
in_shape(d1, [{"text": "✕  REJECTED — FULL AUTHORING ON DAY ONE", "size": 12, "bold": True, "color": RED, "sa": 6},
              {"text": "Maximum flexibility, but the taxonomy was mid-consolidation, approvals weren't "
               "established, and question changes cascade into compliance reports, dashboards, and "
               "subscription logic.", "size": 13, "color": MUTED, "ls": 1.3}])
d2 = rect(s, 6.72, 3.05, 5.78, 1.85, fill=SURFACE, line=LINE)
in_shape(d2, [{"text": "✓  CHOSEN — CONFIGURATION FIRST, AUTHORING LATER", "size": 12, "bold": True, "color": LEAF, "sa": 6},
              {"text": "Phase 1 introduces real operational control while deferring authoring until "
               "governance matures. Content stays stable; only behavior is configurable. Build trust, "
               "then expand.", "size": 13, "color": MUTED, "ls": 1.3}])
p1 = rect(s, 0.78, 5.1, 5.78, 1.55, fill=SURFACE, line=LINE)
rect(s, 0.78, 5.1, 5.78, 0.06, fill=LEAF, radius=0)
in_shape(p1, [{"text": "PHASE 1 · SHIPPED", "size": 12, "bold": True, "color": LEAF, "sa": 5},
              {"text": "Show/hide · reorder · required toggle · country rules · compliance locks · full audit trail",
               "size": 13, "color": MUTED, "ls": 1.3}])
p2 = rect(s, 6.72, 5.1, 5.78, 1.55, fill=SURFACE, line=LINE)
in_shape(p2, [{"text": "UPCOMING PHASES", "size": 12, "bold": True, "color": MUTED, "sa": 5},
              {"text": "Create/edit/remove questions · sub-question logic · category management · "
               "versioning & rollback · AI-assisted config", "size": 13, "color": MUTED, "ls": 1.3}])
footer(s, "Warden"); pagenum(s, "06")
notes(s, "Your standout strategic moment — 'what I chose NOT to do.' Hiring panels remember "
       "trade-offs more than pixels. The line: 'Phase 1 earns trust before expanding capability, and "
       "I made the roadmap visible so phasing read as strategy, not a limitation.'")

# ============================================================ W7 AI
s = slide()
kicker(s, "Scope Beyond the Platform")
text(s, 0.78, 1.05, 11.7, 0.7, [{"text": "Evaluating AI in high-stakes workflows", "size": 30, "bold": True, "color": INK}])
text(s, 0.78, 2.0, 6.2, 0.45, [{"text": "AI-assisted incident reporting", "size": 19, "bold": True, "color": INK}])
text(s, 0.78, 2.55, 6.2, 1.2, [{"runs": [("Store managers describe a situation in plain language; "
     "the AI pre-populates the form as a ", MUTED, False), ("draft", INK, True), (" — every field "
     "editable, ambiguities flagged. The AI never submits on their behalf.", MUTED, False)],
     "size": 14.5, "ls": 1.35}])
text(s, 0.78, 4.0, 6.2, 0.45, [{"text": "AI-assisted incident summary", "size": 19, "bold": True, "color": INK}])
text(s, 0.78, 4.55, 6.2, 1.2, [{"text": "Operators managing concurrent incidents can surface a "
     "one-paragraph summary on demand — enough to assess severity and act. Always on demand, never "
     "auto-shown.", "size": 14.5, "color": MUTED, "ls": 1.35}])
co = rect(s, 7.3, 2.4, 5.2, 3.4, fill=SURFACE, line=LINE)
rect(s, 7.3, 2.4, 0.06, 3.4, fill=LEAF, radius=0)
in_shape(co, [{"text": "DESIGN PRINCIPLE", "size": 11, "bold": True, "color": MUTED, "sa": 8},
              {"text": "AI assists documentation. Humans remain accountable.", "size": 19, "bold": True, "color": INK, "sa": 12, "ls": 1.2},
              {"text": "The risk in a high-stakes workflow isn't AI extracting data wrong — it's an "
               "operator trusting output they didn't verify. Every pattern was designed around that.",
               "size": 14, "color": MUTED, "ls": 1.4}], anchor=MSO_ANCHOR.MIDDLE)
footer(s, "Warden"); pagenum(s, "07")
notes(s, "Keep this tight but include it — AI literacy is a 2026 hiring signal. The one line to nail: "
       "'The risk isn't bad extraction, it's an operator trusting output they didn't verify.' That "
       "shows you think about AI in terms of accountability, not features.")

# ============================================================ W8 impact + reflection
s = slide()
kicker(s, "Impact & Reflection")
text(s, 0.78, 1.05, 11.7, 0.7, [{"text": "What was delivered — and what I'd do differently", "size": 30, "bold": True, "color": INK}])
tiles = [("0→1","Super Admin configuration platform — no comparable tooling existed before"),
         ("3","surfaces aligned under a unified governance & design-system strategy for the first time"),
         ("∞→1","design-system layers reduced to a single documented alignment strategy with LD 3.5")]
tw = 3.78
for i,(n,t) in enumerate(tiles):
    cx = 0.78 + i*(tw+0.18)
    cd = rect(s, cx, 2.1, tw, 2.0, fill=SURFACE, line=LINE)
    in_shape(cd, [{"text": n, "size": 40, "bold": True, "color": BLUE, "sa": 8},
                  {"text": t, "size": 13, "color": MUTED, "ls": 1.3}])
text(s, 0.78, 4.5, 5.78, 1.8, [{"runs": [("Honest note: ", INK, True), ("success measures are "
     "directional — there was almost no baseline data. If I did it again, I'd push for "
     "instrumentation from day one.", MUTED, False)], "size": 15, "ls": 1.4}])
text(s, 6.72, 4.5, 5.78, 1.8, [{"runs": [("The most valuable outcome was the ", MUTED, False),
     ("governance framework itself", INK, True), (": a clear boundary between configurable and "
     "engineering-only, an audit trail as accountability, and a phased model that builds trust.",
     MUTED, False)], "size": 15, "ls": 1.4}])
footer(s, "Warden"); pagenum(s, "08")
notes(s, "Be candid the metrics are directional — then immediately convert that into your reflection "
       "(instrument from day one). Self-awareness reads as senior. Close on: 'the governance "
       "framework itself was the real deliverable — it's the infrastructure that makes the platform "
       "improvable over time.' Then bridge to Case 2.")

# ============================================================ DIVIDER VISIT
s = divider("Case 02",
    [("When the welcome experience ", WHITE, True), ("isn't", C("7FD28A"), True), (" welcoming", WHITE, True)],
    "Unifying Walmart campus visitor check-in across every US & India campus · Lead Product Designer",
    C("7FD28A"))
notes(s, "Bridge from Warden: 'That was strategy at the system level. This next one is execution at "
       "service scale — and it ships with real numbers.' Then the hook on the next slide.")

# ============================================================ V1 title
s = slide()
dot_eyebrow(s, "Service Design · Facilities & Access · Accessibility")
text(s, 0.78, 1.1, 8.4, 1.1, [{"runs": [("A $500B company with a ", INK, True), ("clipboard", BLUE, True),
     (" at the front door", INK, True)], "size": 30, "bold": True, "ls": 1.08}])
text(s, 0.78, 2.6, 5.85, 1.4, [{"text": "Twelve campuses, twelve different check-in processes — "
     "fragmented, and none designed with accessibility in mind. I unified them into one accessible "
     "service.", "size": 17, "color": MUTED, "ls": 1.45}])
meta = [("Role","Lead Product Designer"),("Touchpoints","Kiosk · Web · Email · Admin"),
        ("Team","PM · Eng · Security · Ops"),("Timeline","Nov 2024 – Oct 2025")]
mw = 5.85/4
rect(s, 0.78, 4.5, 5.85, 1.2, fill=SURFACE, line=LINE)
for i,(lab,val) in enumerate(meta):
    if i>0: rect(s, 0.78+i*mw, 4.5, 0.012, 1.2, fill=LINE, radius=0)
    text(s, 0.92+i*mw, 4.66, mw-0.2, 1.0, [
        {"text": lab.upper(), "size": 9, "bold": True, "color": MUTED, "sa": 5},
        {"text": val, "size": 12, "bold": True, "color": INK, "ls": 1.15}])
picture(s, "kiosk-real.jpg", 6.95, 1.1, box_w=5.6, box_h=4.6)
rect(s, 6.95, 1.1, 5.6, 4.6, line=LINE)
footer(s, "Visit Check-in"); pagenum(s, "01")
notes(s, "Use the hook line verbatim — 'a $500B company greeting visitors with a clipboard.' It's "
       "memorable and gets a reaction. Then: twelve campuses, twelve systems, none accessible. Set "
       "up that the hardest part had nothing to do with screens.")

# ============================================================ V2 constraints
s = slide()
kicker(s, "My Role & the constraints")
text(s, 0.78, 1.05, 11, 0.7, [{"text": "Sole lead designer, end to end", "size": 30, "bold": True, "color": INK}])
text(s, 0.78, 1.85, 6.4, 1.1, [{"text": "I owned every surface — kiosk, preregistration web flow, "
     "email system, and admin portal — and led the accessibility strategy. WCAG 2.1 AA shaped the "
     "interaction model from day one.", "size": 15, "color": MUTED, "ls": 1.4}])
ccards = [("Hardware","Surface Go's compact screen — every interaction one decision deep."),
          ("Connectivity","Unreliable at several sites — full offline mode with paper fallbacks."),
          ("Inclusivity","Motion sensitivity, mount heights, screen-reader support — designed in."),
          ("Fragmentation","Twelve implementations, no shared library, no cross-campus visibility.")]
gw2, gh2 = 3.12, 1.3
for i,(tt,bb) in enumerate(ccards):
    cx = 0.78 + (i%2)*(gw2+0.16); cy = 3.3 + (i//2)*(gh2+0.16)
    cd = rect(s, cx, cy, gw2, gh2, fill=SURFACE, line=LINE)
    in_shape(cd, [{"text": tt.upper(), "size": 11.5, "bold": True, "color": INK, "sa": 5},
                  {"text": bb, "size": 12, "color": MUTED, "ls": 1.25}])
picture(s, "field-wheelchair.jpg", 7.5, 2.0, box_w=5.0, box_h=3.0)
rect(s, 7.5, 2.0, 5.0, 3.0, line=LINE)
text(s, 7.5, 5.15, 5.0, 1.4, [{"runs": [("Accessibility gap: ", INK, True), ("for some visitors, "
     "check-in simply didn't work — no keyboard path, no contrast compliance. A non-negotiable "
     "constraint from day one.", MUTED, False)], "size": 12.5, "ls": 1.3}])
footer(s, "Visit Check-in"); pagenum(s, "02")
notes(s, "Lead with the wheelchair photo — it's your emotional anchor and makes accessibility "
       "concrete, not a checkbox. Mention a security associate had flagged this gap for over a year. "
       "Then the constraints grid shows you designed within real hardware/connectivity limits.")

# ============================================================ V3 discovery
s = slide()
kicker(s, "Discovery")
text(s, 0.78, 1.05, 11, 0.7, [{"text": "The problem the software couldn't see", "size": 30, "bold": True, "color": INK}])
text(s, 0.78, 1.85, 11.7, 0.85, [{"runs": [("Every implementation was built inside-out — what does "
     "the system need to capture? Nobody asked what the visitor, host, or associate needed to ", MUTED, False),
     ("do", INK, True), (". The screens were fine; the service around them had no connective tissue.",
     MUTED, False)], "size": 16, "ls": 1.4}])
# pull quote
rect(s, 0.78, 2.95, 0.06, 1.0, fill=BLUE, radius=0)
text(s, 1.0, 2.95, 11.4, 1.0, [
    {"text": "“I end up re-entering half the visitors myself because the screen times out on "
     "them. None of that shows up anywhere.”", "size": 20, "bold": True, "color": INK, "sa": 6, "ls": 1.25},
    {"text": "— Front-desk associate, field interview", "size": 13, "color": MUTED}])
findings = [("1 · Preregistration","The highest-leverage fix nobody had prioritized — the tool existed in a separate system nobody knew about."),
            ("2 · Accessibility","Failures created invisible operational load — associates completing check-in for visitors who couldn't. The workaround became the process."),
            ("3 · Host awareness","Hosts had no real-time signal their visitor arrived. Visitors waited in the lobby; no system solved it.")]
fw = 3.78
for i,(tt,bb) in enumerate(findings):
    cx = 0.78 + i*(fw+0.18)
    cd = rect(s, cx, 4.5, fw, 2.0, fill=SURFACE, line=LINE)
    in_shape(cd, [{"text": tt, "size": 13, "bold": True, "color": INK, "sa": 6},
                  {"text": bb, "size": 12.5, "color": MUTED, "ls": 1.3}])
footer(s, "Visit Check-in"); pagenum(s, "03")
notes(s, "The reframe for this case: systems were built inside-out (what to capture) instead of "
       "around what people need to do. Read the associate quote out loud — it proves the workaround "
       "had become the invisible process. Three findings changed the project's direction.")

# ============================================================ V4 alignment
s = slide()
kicker(s, "The hard part")
text(s, 0.78, 1.05, 11.7, 0.7, [{"text": "Getting everyone to agree was the real design challenge",
     "size": 28, "bold": True, "color": INK}])
text(s, 0.78, 2.1, 6.2, 1.5, [{"runs": [("Each campus had a different definition of 'welcoming,' and "
     "sequential reviews stalled progress. Leadership engaged with visuals more readily than service "
     "flows — so ", MUTED, False), ("the landing page became a proxy for confidence in the entire "
     "system.", INK, True)], "size": 16, "ls": 1.4}])
text(s, 0.78, 4.0, 6.2, 1.5, [{"text": "I recognized what was happening: stakeholders hadn't built a "
     "shared mental model, so the landing page became where anxiety surfaced.", "size": 16, "color": MUTED, "ls": 1.4}])
co = rect(s, 7.3, 2.1, 5.2, 3.6, fill=SURFACE, line=LINE)
rect(s, 7.3, 2.1, 0.06, 3.6, fill=BLUE, radius=0)
in_shape(co, [{"text": "WHAT I CHANGED", "size": 11, "bold": True, "color": MUTED, "sa": 8},
              {"runs": [("Instead of sequential reviews, I presented ", INK, False), ("all directions "
               "simultaneously", BLUEINK, True), (" — video, carousel, static — each with explicit "
               "tradeoffs tied to goals we'd already aligned on. Making the decision space visible let "
               "stakeholders converge without feeling overruled.", INK, False)], "size": 15.5, "ls": 1.45}],
         anchor=MSO_ANCHOR.MIDDLE)
footer(s, "Visit Check-in"); pagenum(s, "04")
notes(s, "This is your design-LEADERSHIP slide, not craft. The insight: the landing-page debate was "
       "really a missing shared mental model. You diagnosed the meeting dynamic, not just the UI, and "
       "restructured the process. This is what separates senior from mid-level — emphasize it.")

# ============================================================ V5 design
s = slide()
kicker(s, "The Design")
text(s, 0.78, 1.05, 11.7, 0.7, [{"text": "Unified UI, local logic — one system, different campuses",
     "size": 28, "bold": True, "color": INK}])
picture(s, "journey-map-walkin.png", 0.78, 2.0, box_w=7.2, box_h=4.1)
rect(s, 0.78, 2.0, 7.2, 4.1, line=LINE)
text(s, 0.78, 6.2, 7.2, 0.5, [{"runs": [("Walk-in visitor journey: ", INK, True), ("the opportunities "
     "row drove every kiosk screen — each traces to a field observation.", MUTED, False)], "size": 12, "ls": 1.25}])
text(s, 8.3, 2.0, 4.25, 0.9, [{"runs": [("The answer was a ", MUTED, False), ("configurable flow "
     "architecture", INK, True), (": the same interface, with location-aware steps that activate by "
     "facility type.", MUTED, False)], "size": 15, "ls": 1.35}])
flows = ["NDA signing for California","Device registration for India (IDC)","Streamlined short-form for 8th & Plate"]
for i,fl in enumerate(flows):
    ty = 3.65 + i*0.5
    rect(s, 8.3, ty+0.04, 0.16, 0.16, fill=LEAF, radius=0.5)
    text(s, 8.58, ty, 4.0, 0.45, [{"text": fl, "size": 14, "color": MUTED}])
text(s, 8.3, 5.3, 4.25, 1.2, [{"runs": [("All on one codebase, one Surface Go. Walk-in kiosk operable "
     "in ", MUTED, False), ("under 60 seconds without instruction", INK, True), (" — progressive "
     "disclosure, one decision per screen.", MUTED, False)], "size": 14, "ls": 1.3}])
footer(s, "Visit Check-in"); pagenum(s, "05")
notes(s, "The systems-thinking payoff: one configurable flow, not twelve forks. Location-aware steps "
       "activate by facility type on one codebase. Tie it back to the journey map — every screen "
       "decision traces to a field observation, so it's research-driven, not aesthetic.")

# ============================================================ V6 loops
s = slide()
kicker(s, "Closing the loops")
text(s, 0.78, 1.05, 11.7, 0.7, [{"text": "Three capabilities the legacy system never had", "size": 30, "bold": True, "color": INK}])
loops = [("Check-out & multi-building","Visitors scan at each entry/exit — no repeated check-ins. Every scan syncs in real time to the Visit Admin portal: a complete, time-stamped record of campus movement. Net-new visibility."),
         ("Preregistration","The tool existed but nobody knew. I wired the prereg link directly into the host's calendar invite — so sending the invite and the prep link became a single action. Pre-registered visitors clear far faster."),
         ("Host notification","Any completed check-in triggers an instant alert to the host — defaulting to no action needed. One tap to decline; the front desk is then notified automatically. No phone call, no manual escalation.")]
lw = 3.78
for i,(tt,bb) in enumerate(loops):
    cx = 0.78 + i*(lw+0.18)
    cd = rect(s, cx, 2.4, lw, 3.7, fill=SURFACE, line=LINE)
    in_shape(cd, [{"text": tt.upper(), "size": 13, "bold": True, "color": INK, "sa": 8},
                  {"text": bb, "size": 13, "color": MUTED, "ls": 1.4}])
footer(s, "Visit Check-in"); pagenum(s, "06")
notes(s, "Three net-new loops the legacy system never closed. The prereg story is the strongest — "
       "the highest-leverage fix was organizational, not visual: the tool existed but was invisible, "
       "so you made sending it a single action with the calendar invite. That's service design.")

# ============================================================ V7 accessibility
s = slide()
kicker(s, "Accessibility built into the practice")
text(s, 0.78, 1.05, 11.7, 0.7, [{"text": "Keyboard-first, validated before handoff", "size": 30, "bold": True, "color": INK}])
steps = [("1","Built the interaction model keyboard-navigation first — every touch interaction had a keyboard equivalent from day one."),
         ("2","Every screen validated against a shared accessibility checklist before it went to engineering."),
         ("3","Annotations embedded directly in Figma — contrast ratios, focus order, screen-reader labels — so developers had zero ambiguity at handoff.")]
for i,(n,b) in enumerate(steps):
    ty = 2.2 + i*1.15
    rect(s, 0.78, ty, 0.42, 0.42, fill=BLUEPALE, radius=0.5)
    text(s, 0.78, ty+0.07, 0.42, 0.4, [{"text": n, "size": 14, "bold": True, "color": BLUEINK, "align": PP_ALIGN.CENTER}])
    text(s, 1.4, ty, 5.3, 1.0, [{"text": b, "size": 15, "color": MUTED, "ls": 1.35}])
co = rect(s, 7.3, 2.4, 5.2, 2.9, fill=SURFACE, line=LINE)
rect(s, 7.3, 2.4, 0.06, 2.9, fill=LEAF, radius=0)
in_shape(co, [{"text": "WHY IT MATTERS", "size": 11, "bold": True, "color": MUTED, "sa": 8},
              {"runs": [("A wheelchair user at the Bentonville turnstile had no guaranteed accessible "
               "check-in path. A security associate had flagged this gap for ", INK, False),
               ("over a year", BLUEINK, True), (" before design started.", INK, False)], "size": 16, "ls": 1.45}],
         anchor=MSO_ANCHOR.MIDDLE)
footer(s, "Visit Check-in"); pagenum(s, "07")
notes(s, "Concrete craft evidence: keyboard-first model, shared checklist, annotations embedded in "
       "the Figma handoff so devs had zero ambiguity. The 'flagged for over a year' detail shows the "
       "human cost of the gap you closed. Keep it brief — this is proof, not the headline.")

# ============================================================ V8 impact
s = slide()
kicker(s, "Outcomes")
text(s, 0.78, 1.05, 11.7, 0.7, [{"text": "What changed, and how we measured it", "size": 30, "bold": True, "color": INK}])
text(s, 0.78, 1.85, 11.7, 0.6, [{"runs": [("Shipped to all US campuses, then Bengaluru and other "
     "India locations — becoming the ", MUTED, False), ("global standard", INK, True), (" for Walmart "
     "campus visitor management.", MUTED, False)], "size": 16, "ls": 1.4}])
tiles = [("30%","reduction in average visitor processing time across all campuses post-launch"),
         ("3×","increase in preregistration completion rate"),
         ("~60%","fewer front-desk escalations at high-traffic campuses in the first quarter")]
tw = 3.78
for i,(n,t) in enumerate(tiles):
    cx = 0.78 + i*(tw+0.18)
    cd = rect(s, cx, 2.85, tw, 2.0, fill=SURFACE, line=LINE)
    in_shape(cd, [{"text": n, "size": 44, "bold": True, "color": BLUE, "sa": 8},
                  {"text": t, "size": 13, "color": MUTED, "ls": 1.3}])
text(s, 0.78, 5.2, 11.7, 1.2, [{"text": "The admin portal eliminated the most common manual "
     "override — associates re-entering visitor data after kiosk timeouts — which had driven a "
     "significant share of front-desk escalations.", "size": 14.5, "color": MUTED, "ls": 1.4}])
footer(s, "Visit Check-in"); pagenum(s, "08")
notes(s, "Your strongest numbers — lead with them out loud: 30% faster, 3× preregistration, ~60% "
       "fewer escalations, now the global standard. Then the qualitative: eliminating the kiosk-"
       "timeout re-entry workaround is what drove the escalation drop. Numbers first, story second.")

# ============================================================ V9 admin backbone
s = slide()
kicker(s, "The operational backbone · Visit Admin")
text(s, 0.78, 1.05, 11.7, 0.7, [{"text": "Behind the experience is the tool security runs on", "size": 28, "bold": True, "color": INK}])
# before
rect(s, 0.78, 2.0, 1.3, 0.32, fill=REDPALE, radius=0.3)
text(s, 0.78, 2.04, 1.3, 0.3, [{"text": "BEFORE", "size": 10, "bold": True, "color": RED, "align": PP_ALIGN.CENTER}])
picture(s, "admin-wl-b1.png", 0.78, 2.45, box_w=3.5, box_h=2.5)
rect(s, 0.78, 2.45, 3.5, 2.5, line=LINE)
# after
rect(s, 4.5, 2.0, 1.3, 0.32, fill=GREENPALE, radius=0.3)
text(s, 4.5, 2.04, 1.3, 0.3, [{"text": "PHASE 1", "size": 10, "bold": True, "color": GREENINK, "align": PP_ALIGN.CENTER}])
picture(s, "admin-wl-a1.png", 4.5, 2.45, box_w=3.5, box_h=2.5)
rect(s, 4.5, 2.45, 3.5, 2.5, line=LINE)
text(s, 0.78, 5.1, 7.2, 0.9, [{"runs": [("Watchlist: ", INK, True), ("icon-only threat flags → "
     "semantic labeled badges (Caution / Critical / Awareness). In a safety-critical module, "
     "ambiguity has direct consequences.", MUTED, False)], "size": 13, "ls": 1.3}])
text(s, 8.3, 2.2, 4.25, 1.6, [{"runs": [("When engineering capacity shifted to the kiosk rebuild, I "
     "scoped Phase 1 to ", MUTED, False), ("LD 3.5 migration + targeted usability fixes", INK, True),
     (" and kept the IA work documented and queued.", MUTED, False)], "size": 14.5, "ls": 1.35}])
text(s, 8.3, 4.3, 4.25, 1.8, [{"runs": [("Phasing also de-risked adoption: many associates are ",
     MUTED, False), ("20+ year", INK, True), (" employees on the legacy manual process. Incremental "
     "releases let phasing read as ", MUTED, False), ("strategy, not compromise", INK, True),
     (".", MUTED, False)], "size": 14.5, "ls": 1.35}])
footer(s, "Visit Admin"); pagenum(s, "09")
notes(s, "Fold Admin in as the backbone — two beats only: (1) the Watchlist before/after is the most "
       "vivid (icon-only → labeled, safety-critical), (2) strategic scoping when capacity shifted, "
       "kept IA queued. The 20-year-tenured associate detail makes phasing feel human, not just "
       "process. Don't make this a full third case unless asked.")

# ============================================================ V10 reflection
s = slide()
kicker(s, "Reflection")
text(s, 0.78, 1.05, 11.7, 0.7, [{"text": "What I'd do differently", "size": 30, "bold": True, "color": INK}])
refl = [("1","I'd instrument the service earlier. We had almost no baseline data — check-in times, prereg rates, accessibility escalations were all unmeasured. Measurement should have been scoped in from the start, not advocated for mid-project."),
        ("2","I'd push harder for pilot-then-scale. We launched across all US campuses nearly simultaneously. A single high-traffic pilot with a six-week learning period would have surfaced group-booking edge cases — and saved two engineering sprints."),
        ("3","I underestimated the front-desk associate as a design partner. I interviewed them early but treated them as research subjects. I should have brought two in as co-reviewers — they caught workflow gaps I was too close to see.")]
for i,(n,b) in enumerate(refl):
    ty = 2.0 + i*1.5
    rect(s, 0.78, ty, 0.42, 0.42, fill=BLUEPALE, radius=0.5)
    text(s, 0.78, ty+0.07, 0.42, 0.4, [{"text": n, "size": 14, "bold": True, "color": BLUEINK, "align": PP_ALIGN.CENTER}])
    runs = [(b.split(".")[0]+".", INK, True), (b[len(b.split(".")[0])+1:], MUTED, False)]
    text(s, 1.4, ty, 11.1, 1.3, [{"runs": runs, "size": 15, "ls": 1.35}])
footer(s, "Visit Check-in"); pagenum(s, "10")
notes(s, "Three mature, specific reflections — not generic. Each names a real trade-off: "
       "instrumentation, pilot-then-scale, and treating associates as co-designers. Showing you'd "
       "change your own process reads as senior self-awareness. End here, then go to the close.")

# ============================================================ CLOSING
s = slide(INK)
text(s, 0, 1.5, SW, 0.4, [{"text": "THANK YOU", "size": 13, "bold": True, "color": LT2, "align": PP_ALIGN.CENTER}])
text(s, 1, 2.0, SW-2, 1.4, [{"runs": [("From designing screens to ", WHITE, True),
     ("designing systems", SKY, True)], "size": 40, "align": PP_ALIGN.CENTER, "ls": 1.05}])
b1 = rect(s, 2.9, 3.9, 3.4, 1.6, fill=None)
text(s, 2.9, 4.0, 3.4, 1.5, [{"text": "Warden", "size": 14, "color": LT2, "align": PP_ALIGN.CENTER, "sa": 6},
     {"text": "A governance framework that makes the platform improvable over time.", "size": 14,
      "color": LT, "align": PP_ALIGN.CENTER, "ls": 1.4}])
text(s, 7.0, 4.0, 3.4, 1.5, [{"text": "Visit Check-in", "size": 14, "color": LT2, "align": PP_ALIGN.CENTER, "sa": 6},
     {"text": "A unified, accessible service that became Walmart's global standard.", "size": 14,
      "color": LT, "align": PP_ALIGN.CENTER, "ls": 1.4}])
text(s, 0, 5.85, SW, 0.4, [{"text": "Stanlyn Lu · Product Designer", "size": 15, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}])
text(s, 0, 6.35, SW, 0.4, [{"text": "lintanxi@gmail.com · linkedin.com/in/lintanxi", "size": 13, "color": LT2, "align": PP_ALIGN.CENTER}])
notes(s, "Close with the one-line throughline: 'Both cases are really the same move — looking past "
       "the stated request to the system underneath, then designing the structure that makes things "
       "better over time.' Then invite questions.")

out = "/Users/lintanxi/Downloads/Stanlyn_Lu_Case_Studies.pptx"
prs.save(out)

# --- Keynote compatibility fix ---
# python-pptx adds a notesMaster relationship but omits the <p:notesMasterIdLst>
# element in presentation.xml. PowerPoint tolerates this; Keynote rejects the whole
# file ("file format is invalid"). Inject the missing element so Keynote accepts it.
import zipfile, shutil, re
def fix_notes(path):
    tmp = path + ".tmp"
    zin = zipfile.ZipFile(path, "r")
    rels = zin.read("ppt/_rels/presentation.xml.rels").decode("utf-8")
    m = (re.search(r'<Relationship[^>]*Type="[^"]*/notesMaster"[^>]*Id="(rId\d+)"', rels) or
         re.search(r'<Relationship[^>]*Id="(rId\d+)"[^>]*Type="[^"]*/notesMaster"', rels))
    if not m:
        zin.close(); return
    rid = m.group(1)
    pres = zin.read("ppt/presentation.xml").decode("utf-8")
    if "notesMasterIdLst" not in pres:
        inject = '<p:notesMasterIdLst><p:notesMasterId r:id="%s"/></p:notesMasterIdLst>' % rid
        pres = pres.replace("</p:sldMasterIdLst>", "</p:sldMasterIdLst>" + inject, 1)
    zout = zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED)
    for it in zin.infolist():
        d = pres.encode("utf-8") if it.filename == "ppt/presentation.xml" else zin.read(it.filename)
        zout.writestr(it, d)
    zout.close(); zin.close(); shutil.move(tmp, path)

fix_notes(out)
print("Saved:", out, "slides:", len(prs.slides._sldIdLst), "(Keynote-compatible)")
